from __future__ import annotations

import csv
import io
import json
import mimetypes
import re
import zipfile
from dataclasses import dataclass
from pathlib import PurePosixPath
from typing import Iterable

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from ukb.ingestion_models import (
    IngestionItemStatus,
    IngestionPreview,
    IngestionPreviewItem,
    IngestionSourceMode,
)

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".json",
    ".yaml",
    ".yml",
    ".sql",
    ".html",
    ".htm",
    ".xml",
    ".rst",
    ".log",
}
BINARY_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | BINARY_EXTENSIONS
ARCHIVE_EXTENSION = ".zip"


@dataclass(slots=True)
class RawIngestionItem:
    name: str
    path: str
    data: bytes
    content_type: str = "application/octet-stream"
    source_uri: str | None = None


@dataclass(slots=True)
class ParsedIngestionItem:
    raw: RawIngestionItem
    text: str
    status: IngestionItemStatus
    warnings: list[str]


class IngestionParserService:
    """Parse bounded source batches and produce a reviewable manifest.

    The parser is deliberately deterministic. LLM enrichment happens only after
    this layer has established the source, path, MIME type and normalized text.
    """

    def __init__(
        self,
        *,
        max_file_bytes: int,
        max_batch_files: int,
        max_archive_bytes: int,
        max_archive_uncompressed_bytes: int,
    ) -> None:
        self.max_file_bytes = max_file_bytes
        self.max_batch_files = max_batch_files
        self.max_archive_bytes = max_archive_bytes
        self.max_archive_uncompressed_bytes = max_archive_uncompressed_bytes

    def preview(
        self,
        items: list[RawIngestionItem],
        *,
        source_mode: IngestionSourceMode,
        connector: str,
    ) -> tuple[IngestionPreview, list[ParsedIngestionItem]]:
        expanded = self._expand_archives(items) if source_mode == IngestionSourceMode.zip else items
        if len(expanded) > self.max_batch_files:
            expanded = expanded[: self.max_batch_files]
            batch_warning = (
                f"The batch exceeded {self.max_batch_files} files; only the first "
                f"{self.max_batch_files} were inspected."
            )
        else:
            batch_warning = None

        parsed = [self.parse_item(item) for item in expanded]
        preview_items = [
            IngestionPreviewItem(
                name=item.raw.name,
                path=item.raw.path,
                content_type=item.raw.content_type,
                size_bytes=len(item.raw.data),
                status=item.status,
                extracted_chars=len(item.text),
                source_uri=item.raw.source_uri,
            )
            for item in parsed
        ]
        warnings = [warning for item in parsed for warning in item.warnings]
        if batch_warning:
            warnings.insert(0, batch_warning)
        rejected = [item.raw.path for item in parsed if item.status == IngestionItemStatus.rejected]
        usable = [
            item
            for item in parsed
            if item.status != IngestionItemStatus.rejected and item.text.strip()
        ]
        preview = IngestionPreview(
            source_mode=source_mode,
            ready=bool(usable),
            items=preview_items,
            warnings=warnings,
            rejected_items=rejected,
            extracted_chars=sum(len(item.text) for item in usable),
            preview_markdown=self._preview_markdown(usable),
            connector=connector,
        )
        return preview, parsed

    def parse_item(self, item: RawIngestionItem) -> ParsedIngestionItem:
        path = self._safe_path(item.path)
        if not item.data:
            return ParsedIngestionItem(
                raw=item,
                text="",
                status=IngestionItemStatus.rejected,
                warnings=[f"{path}: empty source rejected."],
            )
        if len(item.data) > self.max_file_bytes:
            return ParsedIngestionItem(
                raw=item,
                text="",
                status=IngestionItemStatus.rejected,
                warnings=[
                    f"{path}: {len(item.data)} bytes exceeds the "
                    f"{self.max_file_bytes}-byte file limit."
                ],
            )

        suffix = PurePosixPath(path).suffix.lower()
        if suffix not in SUPPORTED_EXTENSIONS:
            return ParsedIngestionItem(
                raw=item,
                text="",
                status=IngestionItemStatus.rejected,
                warnings=[f"{path}: unsupported format {suffix or '(no extension)'}."],
            )

        try:
            text = self._extract(item.data, suffix)
        except Exception as exc:
            return ParsedIngestionItem(
                raw=item,
                text="",
                status=IngestionItemStatus.rejected,
                warnings=[f"{path}: parser failed ({type(exc).__name__}: {exc})."],
            )

        normalized = self._normalize(text)
        if not normalized:
            return ParsedIngestionItem(
                raw=item,
                text="",
                status=IngestionItemStatus.rejected,
                warnings=[f"{path}: parser produced no usable text."],
            )

        item_warnings: list[str] = []
        status = IngestionItemStatus.ready
        if len(normalized) < 80:
            status = IngestionItemStatus.warning
            item_warnings.append(f"{path}: extracted text is very short; confirm completeness.")
        return ParsedIngestionItem(raw=item, text=normalized, status=status, warnings=item_warnings)

    def _extract(self, data: bytes, suffix: str) -> str:
        if suffix in {".txt", ".md", ".markdown", ".sql", ".rst", ".log", ".yaml", ".yml"}:
            return data.decode("utf-8-sig")
        if suffix == ".csv":
            return self._csv_to_markdown(data.decode("utf-8-sig"))
        if suffix == ".json":
            parsed = json.loads(data.decode("utf-8-sig"))
            return json.dumps(parsed, indent=2, ensure_ascii=False)
        if suffix in {".html", ".htm", ".xml"}:
            soup = BeautifulSoup(data, "html.parser" if suffix != ".xml" else "xml")
            for unwanted in soup(["script", "style", "noscript", "template", "svg"]):
                unwanted.decompose()
            title = soup.title.get_text(" ", strip=True) if soup.title else ""
            body = soup.get_text("\n", strip=True)
            return f"# {title}\n\n{body}" if title else body
        if suffix == ".pdf":
            reader = PdfReader(io.BytesIO(data))
            return "\n\n".join(
                f"## Page {index + 1}\n\n{page.extract_text() or ''}"
                for index, page in enumerate(reader.pages)
            )
        if suffix == ".docx":
            document = Document(io.BytesIO(data))
            blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
            for table_index, table in enumerate(document.tables, start=1):
                blocks.append(f"\n## Table {table_index}\n")
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                blocks.append(self._rows_to_markdown(rows))
            return "\n\n".join(blocks)
        if suffix == ".pptx":
            presentation = Presentation(io.BytesIO(data))
            slides: list[str] = []
            for index, slide in enumerate(presentation.slides, start=1):
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                slides.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
            return "\n\n".join(slides)
        if suffix == ".xlsx":
            workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            sheets: list[str] = []
            for sheet in workbook.worksheets:
                rows = [
                    ["" if value is None else str(value) for value in row]
                    for row in sheet.iter_rows(values_only=True)
                ]
                rows = [row for row in rows if any(cell.strip() for cell in row)]
                sheets.append(f"## Sheet: {sheet.title}\n\n{self._rows_to_markdown(rows)}")
            return "\n\n".join(sheets)
        raise ValueError(f"No parser registered for {suffix}")

    def _expand_archives(self, items: list[RawIngestionItem]) -> list[RawIngestionItem]:
        if len(items) != 1:
            raise ValueError("ZIP ingestion accepts exactly one archive per batch.")
        archive = items[0]
        if len(archive.data) > self.max_archive_bytes:
            raise ValueError(
                f"Archive exceeds the {self.max_archive_bytes}-byte compressed size limit."
            )
        extracted: list[RawIngestionItem] = []
        total = 0
        with zipfile.ZipFile(io.BytesIO(archive.data)) as package:
            for entry in package.infolist():
                if entry.is_dir():
                    continue
                safe_path = self._safe_path(entry.filename)
                if entry.flag_bits & 0x1:
                    raise ValueError(f"Encrypted archive entry is not supported: {safe_path}")
                total += entry.file_size
                if total > self.max_archive_uncompressed_bytes:
                    raise ValueError(
                        "Archive expands beyond the configured uncompressed-size limit."
                    )
                data = package.read(entry)
                content_type = mimetypes.guess_type(safe_path)[0] or "application/octet-stream"
                extracted.append(
                    RawIngestionItem(
                        name=PurePosixPath(safe_path).name,
                        path=safe_path,
                        data=data,
                        content_type=content_type,
                        source_uri=f"zip://{archive.name}/{safe_path}",
                    )
                )
                if len(extracted) > self.max_batch_files:
                    break
        return extracted

    @staticmethod
    def _safe_path(path: str) -> str:
        raw = path.replace("\\", "/")
        if raw.startswith("/") or re.match(r"^[A-Za-z]:", raw):
            raise ValueError(f"Absolute source path is not allowed: {path!r}")
        parts = PurePosixPath(raw).parts
        if not raw or any(part in {"", ".", ".."} for part in parts):
            raise ValueError(f"Unsafe source path: {path!r}")
        return str(PurePosixPath(*parts))

    @staticmethod
    def _normalize(text: str) -> str:
        text = text.replace("\x00", "")
        lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
        result: list[str] = []
        blank = False
        for line in lines:
            if line:
                result.append(line)
                blank = False
            elif not blank and result:
                result.append("")
                blank = True
        return "\n".join(result).strip()

    @staticmethod
    def _csv_to_markdown(text: str) -> str:
        rows = list(csv.reader(io.StringIO(text)))
        return IngestionParserService._rows_to_markdown(rows)

    @staticmethod
    def _rows_to_markdown(rows: Iterable[Iterable[str]]) -> str:
        materialized = [list(row) for row in rows]
        if not materialized:
            return ""
        width = max(len(row) for row in materialized)
        padded = [row + [""] * (width - len(row)) for row in materialized]
        header = padded[0]
        body = padded[1:]
        lines = [
            "| " + " | ".join(cell.replace("|", "\\|") for cell in header) + " |",
            "| " + " | ".join("---" for _ in header) + " |",
        ]
        lines.extend(
            "| " + " | ".join(cell.replace("|", "\\|") for cell in row) + " |"
            for row in body
        )
        return "\n".join(lines)

    @staticmethod
    def _preview_markdown(items: list[ParsedIngestionItem]) -> str:
        sections: list[str] = []
        for item in items[:3]:
            sections.append(f"# {item.raw.name}\n\n{item.text[:3500]}")
        if len(items) > 3:
            sections.append(f"\n_+ {len(items) - 3} additional parsed sources_\n")
        return "\n\n---\n\n".join(sections)
